"""
Tests for document extraction functionality.

This module validates the document extraction pipeline that converts
uploaded files (PDF, DOCX, XLSX, PPTX, text, code) into structured
content parts (text + images) for the LLM conversation.

Tests cover:
  - Text file extraction (plain text, code files with syntax context)
  - PDF extraction (text pages, scanned page detection, embedded images)
  - DOCX extraction (paragraphs, tables, inline images)
  - XLSX extraction (sheet data as markdown tables, multi-sheet)
  - PPTX extraction (slide text, embedded images)
  - Image compression (resize, format conversion)
  - File size validation (reject files > 25 MB)
  - File type validation (reject unsupported extensions)
  - ContentPart serialization to OpenAI format
  - ExtractionResult metadata (text preview, counts, warnings)
  - API endpoint validation

These tests create actual test files in-memory (not on disk) so they
run fast and don't require external file fixtures.
"""

import io
from unittest.mock import MagicMock

import pytest

from stream.middleware.utils.document_extractor import (
    CODE_EXTENSIONS,
    MAX_DOCUMENT_SIZE,
    MAX_IMAGE_DIMENSION,
    SCANNED_PAGE_CHAR_THRESHOLD,
    SUPPORTED_EXTENSIONS,
    TEXT_EXTENSIONS,
    ContentPart,
    ExtractionResult,
    _compress_image,
    _ext_to_language,
    _extract_text_file,
    _table_to_markdown,
    extract_document,
)

# =============================================================================
# ContentPart Tests
# =============================================================================


class TestContentPart:
    """Tests for the ContentPart data class and its OpenAI format conversion.

    ContentPart is the fundamental unit of extracted content — either text
    or an image. Its to_openai_format() method converts it to the exact
    format that OpenAI's multimodal API expects.
    """

    def test_text_part_to_openai_format(self):
        """Text parts become {"type": "text", "text": "..."} objects."""
        part = ContentPart(type="text", text="Hello, world!")
        result = part.to_openai_format()
        assert result == {"type": "text", "text": "Hello, world!"}

    def test_image_part_to_openai_format(self):
        """Image parts become {"type": "image_url", "image_url": {"url": "data:..."}} objects."""
        part = ContentPart(
            type="image",
            image_base64="abc123",
            image_mime="image/png",
        )
        result = part.to_openai_format()
        assert result == {
            "type": "image_url",
            "image_url": {"url": "data:image/png;base64,abc123"},
        }

    def test_text_part_with_none_text(self):
        """Text parts with None text should return empty string."""
        part = ContentPart(type="text", text=None)
        result = part.to_openai_format()
        assert result == {"type": "text", "text": ""}

    def test_unknown_type_raises_error(self):
        """Unknown content types should raise ValueError."""
        part = ContentPart(type="video")
        with pytest.raises(ValueError, match="Unknown content part type"):
            part.to_openai_format()


# =============================================================================
# ExtractionResult Tests
# =============================================================================


class TestExtractionResult:
    """Tests for ExtractionResult serialization."""

    def test_to_dict_serialization(self):
        """ExtractionResult.to_dict() should produce a JSON-serializable dict."""
        result = ExtractionResult(
            filename="test.txt",
            file_type="txt",
            file_size=100,
            content_parts=[
                ContentPart(type="text", text="Hello"),
                ContentPart(type="image", image_base64="abc", image_mime="image/png"),
            ],
            text_preview="Hello",
            total_text_length=5,
            image_count=1,
            page_count=0,
            warnings=["A warning"],
        )
        d = result.to_dict()
        assert d["filename"] == "test.txt"
        assert d["file_type"] == "txt"
        assert d["file_size"] == 100
        assert len(d["content_parts"]) == 2
        assert d["content_parts"][0]["type"] == "text"
        assert d["content_parts"][1]["type"] == "image"
        assert d["text_preview"] == "Hello"
        assert d["total_text_length"] == 5
        assert d["image_count"] == 1
        assert d["page_count"] == 0
        assert d["warnings"] == ["A warning"]


# =============================================================================
# Text File Extraction Tests
# =============================================================================


class TestTextFileExtraction:
    """Tests for plain text and code file extraction.

    Text files are the simplest case — we read bytes as UTF-8 text.
    Code files get wrapped in markdown code fences for syntax context.
    """

    def test_extract_plain_text(self):
        """Plain text files should be extracted as-is with a document header."""
        content = "Hello, this is a test document.\nWith multiple lines."
        file_bytes = content.encode("utf-8")
        result = _extract_text_file(file_bytes, "readme.txt", ".txt")

        assert result.filename == "readme.txt"
        assert result.file_type == "txt"
        assert len(result.content_parts) == 1
        assert result.content_parts[0].type == "text"
        assert "[Document: readme.txt]" in result.content_parts[0].text
        assert "Hello, this is a test document." in result.content_parts[0].text

    def test_extract_python_code(self):
        """Code files should be wrapped in markdown code fences."""
        code = 'def hello():\n    print("Hello!")\n'
        file_bytes = code.encode("utf-8")
        result = _extract_text_file(file_bytes, "main.py", ".py")

        assert result.file_type == "py"
        assert "```python" in result.content_parts[0].text
        assert 'print("Hello!")' in result.content_parts[0].text
        assert "```" in result.content_parts[0].text

    def test_extract_javascript_code(self):
        """JavaScript files should use 'javascript' language tag."""
        code = "const x = 42;\nconsole.log(x);"
        file_bytes = code.encode("utf-8")
        result = _extract_text_file(file_bytes, "app.js", ".js")
        assert "```javascript" in result.content_parts[0].text

    def test_extract_typescript_code(self):
        """TypeScript files should use 'typescript' language tag."""
        code = "const x: number = 42;"
        file_bytes = code.encode("utf-8")
        result = _extract_text_file(file_bytes, "app.ts", ".ts")
        assert "```typescript" in result.content_parts[0].text

    def test_latin1_fallback(self):
        """Files with invalid UTF-8 should fall back to latin-1 decoding."""
        file_bytes = b"\xff\xfe" + "Hello".encode("latin-1")
        result = _extract_text_file(file_bytes, "data.txt", ".txt")
        assert result.content_parts[0].type == "text"

    def test_csv_file(self):
        """CSV files should be read as plain text (no code fence)."""
        csv = "name,score\nAlice,95\nBob,87"
        file_bytes = csv.encode("utf-8")
        result = _extract_text_file(file_bytes, "data.csv", ".csv")
        assert "```" not in result.content_parts[0].text
        assert "Alice,95" in result.content_parts[0].text


# =============================================================================
# Extension to Language Mapping Tests
# =============================================================================


class TestExtToLanguage:
    """Tests for the file extension to markdown language tag mapping."""

    def test_python_extension(self):
        assert _ext_to_language(".py") == "python"

    def test_javascript_extension(self):
        assert _ext_to_language(".js") == "javascript"

    def test_rust_extension(self):
        assert _ext_to_language(".rs") == "rust"

    def test_unknown_extension(self):
        """Unknown extensions should return empty string (no language tag)."""
        assert _ext_to_language(".xyz") == ""


# =============================================================================
# Table to Markdown Tests
# =============================================================================


class TestTableToMarkdown:
    """Tests for converting document tables to markdown format.

    LLMs read markdown tables very well, so this is the best format
    for tabular data — avoids needing to render tables as images.
    """

    def test_simple_table(self):
        """A simple table should be rendered with header separator."""
        table = MagicMock()
        row1 = MagicMock()
        row1.cells = [MagicMock(text="Name"), MagicMock(text="Score")]
        row2 = MagicMock()
        row2.cells = [MagicMock(text="Alice"), MagicMock(text="95")]
        table.rows = [row1, row2]

        result = _table_to_markdown(table)
        assert "| Name | Score |" in result
        assert "| --- | --- |" in result
        assert "| Alice | 95 |" in result

    def test_empty_table(self):
        """An empty table should return empty string."""
        table = MagicMock()
        table.rows = []
        assert _table_to_markdown(table) == ""

    def test_pipe_in_cell_escaped(self):
        """Pipe characters in cell text should be escaped."""
        table = MagicMock()
        row1 = MagicMock()
        row1.cells = [MagicMock(text="A|B")]
        table.rows = [row1]

        result = _table_to_markdown(table)
        assert "A\\|B" in result


# =============================================================================
# Image Compression Tests
# =============================================================================


class TestImageCompression:
    """Tests for the image compression utility.

    Large embedded images are scaled down and compressed as JPEG to
    save payload space. This uses the same approach as the image
    upload pipeline.
    """

    def test_compress_large_image(self):
        """Images larger than MAX_IMAGE_DIMENSION should be resized."""
        from PIL import Image

        img = Image.new("RGB", (2000, 1500), color="red")
        buffer = io.BytesIO()
        img.save(buffer, format="PNG")
        img_bytes = buffer.getvalue()

        compressed, mime = _compress_image(img_bytes, "image/png")

        # Compressed image should be JPEG and smaller
        assert mime == "image/jpeg"
        assert len(compressed) < len(img_bytes)

        # Verify dimensions were reduced
        result_img = Image.open(io.BytesIO(compressed))
        assert result_img.width <= MAX_IMAGE_DIMENSION
        assert result_img.height <= MAX_IMAGE_DIMENSION

    def test_skip_tiny_images(self):
        """Images smaller than 50x50 (icons, bullets) should be returned as-is."""
        from PIL import Image

        img = Image.new("RGB", (30, 30), color="blue")
        buffer = io.BytesIO()
        img.save(buffer, format="PNG")
        img_bytes = buffer.getvalue()

        compressed, mime = _compress_image(img_bytes, "image/png")
        assert compressed == img_bytes
        assert mime == "image/png"

    def test_rgba_conversion(self):
        """RGBA images should be converted to RGB before JPEG compression."""
        from PIL import Image

        img = Image.new("RGBA", (200, 200), color=(255, 0, 0, 128))
        buffer = io.BytesIO()
        img.save(buffer, format="PNG")
        img_bytes = buffer.getvalue()

        compressed, mime = _compress_image(img_bytes, "image/png")
        assert mime == "image/jpeg"

        result_img = Image.open(io.BytesIO(compressed))
        assert result_img.mode == "RGB"


# =============================================================================
# Main extract_document() Tests
# =============================================================================


class TestExtractDocument:
    """Tests for the main extract_document() entry point.

    This function detects file type by extension and routes to the
    appropriate format-specific extractor.
    """

    @pytest.mark.asyncio
    async def test_extract_text_file(self):
        """Text files should be extracted with a document header."""
        content = "Hello, world!"
        result = await extract_document(content.encode("utf-8"), "hello.txt")

        assert result.filename == "hello.txt"
        assert result.file_type == "txt"
        assert result.total_text_length > 0
        assert result.image_count == 0
        assert "Hello, world!" in result.text_preview

    @pytest.mark.asyncio
    async def test_extract_python_file(self):
        """Python files should be wrapped in code fences."""
        code = 'print("Hello!")'
        result = await extract_document(code.encode("utf-8"), "script.py")

        assert result.file_type == "py"
        assert "```python" in result.text_preview

    @pytest.mark.asyncio
    async def test_file_too_large(self):
        """Files exceeding MAX_DOCUMENT_SIZE should raise ValueError."""
        large_bytes = b"x" * (MAX_DOCUMENT_SIZE + 1)
        with pytest.raises(ValueError, match="File too large"):
            await extract_document(large_bytes, "huge.txt")

    @pytest.mark.asyncio
    async def test_unsupported_extension(self):
        """Unsupported file types should raise ValueError."""
        with pytest.raises(ValueError, match="Unsupported file type"):
            await extract_document(b"data", "video.mp4")

    @pytest.mark.asyncio
    async def test_metadata_generation(self):
        """Extraction should generate correct metadata (preview, counts)."""
        text = "A" * 1000
        result = await extract_document(text.encode("utf-8"), "long.txt")

        assert len(result.text_preview) == 500
        assert result.total_text_length > 500
        assert result.image_count == 0

    @pytest.mark.asyncio
    async def test_json_file(self):
        """JSON files should be extracted as plain text."""
        json_content = '{"key": "value", "count": 42}'
        result = await extract_document(json_content.encode("utf-8"), "data.json")
        assert result.file_type == "json"
        assert '"key": "value"' in result.text_preview


# =============================================================================
# PDF Extraction Tests
# =============================================================================


class TestPDFExtraction:
    """Tests for PDF document extraction.

    These tests create simple PDFs in-memory using PyMuPDF (fitz)
    to verify the extraction pipeline without external fixtures.
    """

    @pytest.mark.asyncio
    async def test_extract_text_pdf(self):
        """A PDF with text should extract the text content."""
        import fitz

        doc = fitz.open()
        page = doc.new_page()
        # Insert enough text to exceed the scanned page detection threshold (50 chars).
        # Short text would be misidentified as a scanned page and rendered as image.
        page.insert_text(
            (72, 72),
            "This is a test PDF page with enough text to be recognized as a text page by the extractor.",
        )
        pdf_bytes = doc.tobytes()
        doc.close()

        result = await extract_document(pdf_bytes, "test.pdf")

        assert result.file_type == "pdf"
        assert result.page_count == 1
        assert "test PDF page" in result.text_preview
        assert result.image_count == 0

    @pytest.mark.asyncio
    async def test_extract_multi_page_pdf(self):
        """A multi-page PDF should have page markers in the output."""
        import fitz

        doc = fitz.open()
        for i in range(3):
            page = doc.new_page()
            page.insert_text(
                (72, 72),
                f"Content on page {i + 1}. This line has enough text to exceed the scanned page threshold for detection purposes.",
            )
        pdf_bytes = doc.tobytes()
        doc.close()

        result = await extract_document(pdf_bytes, "multi.pdf")

        assert result.page_count == 3
        text_parts = [p.text for p in result.content_parts if p.type == "text"]
        all_text = "\n".join(t for t in text_parts if t)
        assert "Page 1" in all_text
        assert "Page 2" in all_text
        assert "Page 3" in all_text

    @pytest.mark.asyncio
    async def test_scanned_page_detection(self):
        """Pages with very little text should be rendered as images."""
        import fitz

        doc = fitz.open()
        page = doc.new_page()
        # Insert minimal text (below the threshold)
        page.insert_text((72, 72), "Hi")
        pdf_bytes = doc.tobytes()
        doc.close()

        result = await extract_document(pdf_bytes, "scanned.pdf")

        # Should have at least one image part (the rendered page)
        image_parts = [p for p in result.content_parts if p.type == "image"]
        assert len(image_parts) > 0
        assert any("scanned" in w.lower() for w in result.warnings)


# =============================================================================
# DOCX Extraction Tests
# =============================================================================


class TestDOCXExtraction:
    """Tests for Word document extraction.

    These tests create DOCX files in-memory using python-docx.
    """

    @pytest.mark.asyncio
    async def test_extract_text_docx(self):
        """A DOCX with paragraphs should extract text in order."""
        from docx import Document

        doc = Document()
        doc.add_paragraph("First paragraph.")
        doc.add_paragraph("Second paragraph.")

        buffer = io.BytesIO()
        doc.save(buffer)

        result = await extract_document(buffer.getvalue(), "test.docx")

        assert result.file_type == "docx"
        assert "First paragraph" in result.text_preview
        assert "Second paragraph" in result.text_preview

    @pytest.mark.asyncio
    async def test_extract_docx_with_table(self):
        """A DOCX with a table should render it as markdown."""
        from docx import Document

        doc = Document()
        doc.add_paragraph("Table below:")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Score"
        table.cell(1, 0).text = "Alice"
        table.cell(1, 1).text = "95"

        buffer = io.BytesIO()
        doc.save(buffer)

        result = await extract_document(buffer.getvalue(), "table.docx")

        text_parts = [p.text for p in result.content_parts if p.type == "text" and p.text]
        all_text = "\n".join(text_parts)
        assert "Name" in all_text
        assert "Alice" in all_text
        assert "95" in all_text


# =============================================================================
# XLSX Extraction Tests
# =============================================================================


class TestXLSXExtraction:
    """Tests for Excel spreadsheet extraction.

    Sheet data is converted to markdown tables, which LLMs read well.
    """

    @pytest.mark.asyncio
    async def test_extract_simple_xlsx(self):
        """An XLSX with data should produce a markdown table."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Scores"
        ws.append(["Name", "Score"])
        ws.append(["Alice", 95])
        ws.append(["Bob", 87])

        buffer = io.BytesIO()
        wb.save(buffer)

        result = await extract_document(buffer.getvalue(), "data.xlsx")

        assert result.file_type == "xlsx"
        text_parts = [p.text for p in result.content_parts if p.type == "text" and p.text]
        all_text = "\n".join(text_parts)
        assert "Name" in all_text
        assert "Alice" in all_text
        assert "95" in all_text
        assert "Scores" in all_text

    @pytest.mark.asyncio
    async def test_extract_multi_sheet_xlsx(self):
        """An XLSX with multiple sheets should include all sheets."""
        from openpyxl import Workbook

        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Sheet1"
        ws1.append(["A", "B"])

        ws2 = wb.create_sheet("Sheet2")
        ws2.append(["C", "D"])

        buffer = io.BytesIO()
        wb.save(buffer)

        result = await extract_document(buffer.getvalue(), "multi.xlsx")

        text_parts = [p.text for p in result.content_parts if p.type == "text" and p.text]
        all_text = "\n".join(text_parts)
        assert "Sheet1" in all_text
        assert "Sheet2" in all_text

    @pytest.mark.asyncio
    async def test_empty_sheet_handling(self):
        """Empty sheets should be noted as empty."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Empty"

        buffer = io.BytesIO()
        wb.save(buffer)

        result = await extract_document(buffer.getvalue(), "empty.xlsx")
        text_parts = [p.text for p in result.content_parts if p.type == "text" and p.text]
        all_text = "\n".join(text_parts)
        assert "empty" in all_text.lower()


# =============================================================================
# PPTX Extraction Tests
# =============================================================================


class TestPPTXExtraction:
    """Tests for PowerPoint presentation extraction."""

    @pytest.mark.asyncio
    async def test_extract_text_pptx(self):
        """A PPTX with text slides should extract text per slide."""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
        title = slide.shapes.title
        title.text = "My Presentation"
        body = slide.placeholders[1]
        body.text = "This is the content."

        buffer = io.BytesIO()
        prs.save(buffer)

        result = await extract_document(buffer.getvalue(), "slides.pptx")

        assert result.file_type == "pptx"
        assert result.page_count == 1
        text_parts = [p.text for p in result.content_parts if p.type == "text" and p.text]
        all_text = "\n".join(text_parts)
        assert "My Presentation" in all_text
        assert "This is the content" in all_text


# =============================================================================
# Supported Extensions Tests
# =============================================================================


class TestSupportedExtensions:
    """Tests for the supported file extensions configuration."""

    def test_text_extensions_are_supported(self):
        """All text extensions should be in SUPPORTED_EXTENSIONS."""
        for ext in TEXT_EXTENSIONS:
            assert ext in SUPPORTED_EXTENSIONS, f"{ext} should be supported"

    def test_code_extensions_are_supported(self):
        """All code extensions should be in SUPPORTED_EXTENSIONS."""
        for ext in CODE_EXTENSIONS:
            assert ext in SUPPORTED_EXTENSIONS, f"{ext} should be supported"

    def test_binary_formats_are_supported(self):
        """Binary document formats should be in SUPPORTED_EXTENSIONS."""
        for ext in [".pdf", ".docx", ".xlsx", ".pptx"]:
            assert ext in SUPPORTED_EXTENSIONS, f"{ext} should be supported"

    def test_common_extensions_present(self):
        """Common extensions users would upload should be supported."""
        common = [".py", ".js", ".ts", ".txt", ".md", ".csv", ".json", ".pdf", ".docx"]
        for ext in common:
            assert ext in SUPPORTED_EXTENSIONS, f"{ext} should be supported"


# =============================================================================
# API Endpoint Tests
# =============================================================================


class TestDocumentsAPI:
    """Tests for the /v1/documents/ API endpoints.

    These verify the FastAPI route configuration and error handling
    without making actual HTTP requests.
    """

    def test_router_exists(self):
        """The documents router should be importable."""
        from stream.middleware.routes.documents import router

        assert router is not None

    def test_extract_endpoint_registered(self):
        """The /documents/extract endpoint should be registered."""
        from stream.middleware.routes.documents import router

        routes = [r.path for r in router.routes]
        assert "/documents/extract" in routes

    def test_supported_formats_endpoint_registered(self):
        """The /documents/supported-formats endpoint should be registered."""
        from stream.middleware.routes.documents import router

        routes = [r.path for r in router.routes]
        assert "/documents/supported-formats" in routes


# =============================================================================
# Constants Validation Tests
# =============================================================================


class TestConstants:
    """Tests for configuration constants."""

    def test_max_document_size(self):
        """Max document size should be 25 MB."""
        assert MAX_DOCUMENT_SIZE == 25 * 1024 * 1024

    def test_max_image_dimension(self):
        """Max image dimension should be 1024 pixels."""
        assert MAX_IMAGE_DIMENSION == 1024

    def test_scanned_page_threshold(self):
        """Scanned page threshold should be reasonable (< 100 chars)."""
        assert SCANNED_PAGE_CHAR_THRESHOLD < 100
        assert SCANNED_PAGE_CHAR_THRESHOLD > 0

    def test_supported_extensions_not_empty(self):
        """Supported extensions set should not be empty."""
        assert len(SUPPORTED_EXTENSIONS) > 20
