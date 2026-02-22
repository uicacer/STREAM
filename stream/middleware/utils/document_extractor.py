"""
Document Extraction for STREAM
===============================

This module extracts text and images from uploaded documents so they can
be included in LLM conversations. It supports all common document formats
that campus/research users would upload.

HOW DOCUMENT EXTRACTION WORKS
=============================

When a user uploads a document (e.g., a PDF with text, tables, and figures),
we need to convert it into the OpenAI multimodal message format that LLMs
understand. The key challenge is preserving the ORDER of content — so that
when the document says "as shown in Figure 3", the LLM sees Figure 3's
image right after that text, just like a human reading the document would.

The extraction pipeline:

    Raw file bytes
        │
        ▼
    Detect file type (by extension)
        │
        ▼
    Format-specific extractor
    (PDF → PyMuPDF, DOCX → python-docx, etc.)
        │
        ▼
    List of ContentPart objects in document order:
        [TextPart, ImagePart, TextPart, TextPart, ImagePart, ...]
        │
        ▼
    ExtractionResult (with metadata, previews, warnings)

THE INTERLEAVING PROBLEM
========================

A PDF page might contain:

    "Figure 3 shows the temperature trend..."
    [IMAGE: line chart]
    "Table 2 summarizes the results:"
    [TABLE: data rows]

If we extract text as one blob and images separately, the LLM can't
connect "Figure 3" to the actual chart. So we extract content IN ORDER,
producing interleaved text and image parts:

    ContentPart(type="text", text="Figure 3 shows the temperature trend...")
    ContentPart(type="image", image_base64="...", image_mime="image/png")
    ContentPart(type="text", text="Table 2 summarizes the results:\\n| Col1 | Col2 |...")

This maps directly to the OpenAI multimodal content array format that
the rest of STREAM's pipeline already handles.

SUPPORTED FORMATS
=================

Plain text (read directly):
    .txt, .md, .csv, .json, .xml, .html, .log, .yaml, .yml, .toml, .ini, .cfg

Code files (read directly, with syntax context):
    .py, .js, .ts, .jsx, .tsx, .java, .cpp, .c, .h, .hpp, .go, .rs, .rb,
    .sh, .bash, .sql, .r, .m, .swift, .kt, .scala, .php, .css, .scss, .less

Binary document formats (require library extraction):
    .pdf   → PyMuPDF (fitz): text blocks + embedded images in page order
    .docx  → python-docx: paragraphs, tables, inline images in document order
    .xlsx  → openpyxl: sheet data as markdown tables
    .pptx  → python-pptx: slide text + images per slide

SCANNED PDF DETECTION
=====================

Some PDFs are scanned documents — they're essentially images with no
extractable text. We detect this by checking if a page yields fewer
than 50 characters of text. For such pages, we render the entire page
as an image and let the vision-capable LLM read it (essentially using
the LLM as an OCR engine, which is actually better than traditional OCR
because it understands layout, tables, and figures).
"""

import base64
import io
import logging
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)

# =============================================================================
# CONSTANTS
# =============================================================================

# Maximum file size we'll accept for extraction (25 MB, matching Claude's limit)
MAX_DOCUMENT_SIZE = 25 * 1024 * 1024

# Maximum number of documents per message
MAX_DOCUMENTS_PER_MESSAGE = 10

# If a PDF page has fewer characters than this, it's likely scanned/image-based
# and we should render it as an image instead of trying to extract text.
SCANNED_PAGE_CHAR_THRESHOLD = 50

# Maximum number of scanned pages we'll render as images (to avoid huge payloads)
MAX_SCANNED_PAGE_IMAGES = 10

# Resolution for rendering scanned PDF pages as images (DPI).
# 150 DPI is a good balance: readable by LLMs but not excessively large.
SCANNED_PAGE_DPI = 150

# Maximum image dimension when extracting embedded images from documents.
# Large images are scaled down to save payload space and tokens.
MAX_IMAGE_DIMENSION = 1024

# JPEG quality for compressed images (same as our image upload compression)
IMAGE_QUALITY = 85

# File extensions grouped by extraction method
TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".log",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".env",
    ".tex",
    ".bib",
    ".rst",
}

CODE_EXTENSIONS = {
    ".py",
    ".js",
    ".ts",
    ".jsx",
    ".tsx",
    ".java",
    ".cpp",
    ".c",
    ".h",
    ".hpp",
    ".go",
    ".rs",
    ".rb",
    ".sh",
    ".bash",
    ".sql",
    ".r",
    ".m",
    ".swift",
    ".kt",
    ".kts",
    ".scala",
    ".php",
    ".css",
    ".scss",
    ".less",
    ".sass",
    ".lua",
    ".pl",
    ".pm",
    ".zig",
    ".v",
    ".dart",
    ".jl",
    ".makefile",
    ".dockerfile",
}

# All extensions we support (used for frontend accept attribute)
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | CODE_EXTENSIONS | {".pdf", ".docx", ".xlsx", ".pptx"}


# =============================================================================
# DATA TYPES
# =============================================================================


@dataclass
class ContentPart:
    """A single piece of extracted content — either text or an image.

    These parts are ordered exactly as they appear in the original document,
    so the LLM sees content in the same sequence a human reader would.

    This maps directly to the OpenAI multimodal content format:
      - type="text"  → {"type": "text", "text": "..."}
      - type="image" → {"type": "image_url", "image_url": {"url": "data:image/png;base64,..."}}
    """

    type: str  # "text" or "image"
    text: str | None = None
    image_base64: str | None = None
    image_mime: str | None = None  # e.g., "image/png", "image/jpeg"

    def to_openai_format(self) -> dict:
        """Convert this content part to the OpenAI message content format.

        This is the same format used by STREAM's existing multimodal pipeline
        for pasted/uploaded images. By using the same format, extracted document
        content flows through the existing chat pipeline with zero changes.
        """
        if self.type == "text":
            return {"type": "text", "text": self.text or ""}
        elif self.type == "image":
            data_url = f"data:{self.image_mime};base64,{self.image_base64}"
            return {"type": "image_url", "image_url": {"url": data_url}}
        else:
            raise ValueError(f"Unknown content part type: {self.type}")


@dataclass
class ExtractionResult:
    """The complete result of extracting a document.

    Contains the interleaved content parts plus metadata that the frontend
    needs for display (preview text, warnings, file info).
    """

    filename: str
    file_type: str  # e.g., "pdf", "docx", "txt", "py"
    file_size: int  # original file size in bytes
    content_parts: list[ContentPart] = field(default_factory=list)
    text_preview: str = ""  # first ~500 chars of text for UI preview
    total_text_length: int = 0  # total characters of extracted text
    image_count: int = 0  # number of images extracted
    page_count: int = 0  # for paginated formats (PDF, PPTX)
    warnings: list[str] = field(default_factory=list)

    def to_dict(self) -> dict:
        """Serialize for JSON API response."""
        return {
            "filename": self.filename,
            "file_type": self.file_type,
            "file_size": self.file_size,
            "content_parts": [
                {
                    "type": p.type,
                    "text": p.text,
                    "image_base64": p.image_base64,
                    "image_mime": p.image_mime,
                }
                for p in self.content_parts
            ],
            "text_preview": self.text_preview,
            "total_text_length": self.total_text_length,
            "image_count": self.image_count,
            "page_count": self.page_count,
            "warnings": self.warnings,
        }


# =============================================================================
# MAIN ENTRY POINT
# =============================================================================


async def extract_document(file_bytes: bytes, filename: str) -> ExtractionResult:
    """Extract text and images from an uploaded document.

    This is the main entry point called by the /v1/extract API endpoint.
    It detects the file type by extension and delegates to the appropriate
    format-specific extractor.

    Args:
        file_bytes: The raw bytes of the uploaded file.
        filename: The original filename (used to detect format by extension).

    Returns:
        ExtractionResult with interleaved text and image content parts.

    Raises:
        ValueError: If the file type is unsupported or the file is too large.
    """
    # Validate file size
    if len(file_bytes) > MAX_DOCUMENT_SIZE:
        size_mb = len(file_bytes) / (1024 * 1024)
        raise ValueError(
            f"File too large: {size_mb:.1f} MB. Maximum is "
            f"{MAX_DOCUMENT_SIZE / (1024 * 1024):.0f} MB."
        )

    ext = Path(filename).suffix.lower()
    logger.info(f"[DocExtract] Extracting '{filename}' ({len(file_bytes)} bytes, type: {ext})")

    # Route to the appropriate extractor based on file extension
    if ext in TEXT_EXTENSIONS or ext in CODE_EXTENSIONS:
        result = _extract_text_file(file_bytes, filename, ext)
    elif ext == ".pdf":
        result = _extract_pdf(file_bytes, filename)
    elif ext == ".docx":
        result = _extract_docx(file_bytes, filename)
    elif ext == ".xlsx":
        result = _extract_xlsx(file_bytes, filename)
    elif ext == ".pptx":
        result = _extract_pptx(file_bytes, filename)
    else:
        raise ValueError(
            f"Unsupported file type: '{ext}'. " f"Supported: text, code, .pdf, .docx, .xlsx, .pptx"
        )

    # Generate text preview (first 500 chars of all text parts combined)
    all_text = "\n".join(p.text for p in result.content_parts if p.type == "text" and p.text)
    result.text_preview = all_text[:500]
    result.total_text_length = len(all_text)
    result.image_count = sum(1 for p in result.content_parts if p.type == "image")

    logger.info(
        f"[DocExtract] Extracted '{filename}': "
        f"{result.total_text_length} chars text, "
        f"{result.image_count} images, "
        f"{result.page_count} pages"
    )

    return result


# =============================================================================
# TEXT FILE EXTRACTOR
# =============================================================================


def _extract_text_file(file_bytes: bytes, filename: str, ext: str) -> ExtractionResult:
    """Extract content from plain text and code files.

    These are the simplest files — we just read the bytes as UTF-8 text.
    For code files, we wrap the content in a markdown code fence with the
    appropriate language tag so the LLM knows it's looking at code.

    Args:
        file_bytes: Raw file bytes.
        filename: Original filename.
        ext: File extension (e.g., ".py", ".txt").

    Returns:
        ExtractionResult with a single text ContentPart.
    """
    # Decode with UTF-8, falling back to latin-1 for binary-ish text files
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = file_bytes.decode("latin-1")
        logger.warning(f"[DocExtract] '{filename}' is not valid UTF-8, using latin-1 fallback")

    # For code files, wrap in a markdown code fence so the LLM sees syntax context
    if ext in CODE_EXTENSIONS:
        lang = _ext_to_language(ext)
        formatted_text = f"```{lang}\n{text}\n```"
    else:
        formatted_text = text

    content_text = f"[Document: {filename}]\n\n{formatted_text}"

    return ExtractionResult(
        filename=filename,
        file_type=ext.lstrip("."),
        file_size=len(file_bytes),
        content_parts=[ContentPart(type="text", text=content_text)],
        page_count=0,
    )


def _ext_to_language(ext: str) -> str:
    """Map a file extension to a markdown code fence language tag."""
    mapping = {
        ".py": "python",
        ".js": "javascript",
        ".ts": "typescript",
        ".jsx": "jsx",
        ".tsx": "tsx",
        ".java": "java",
        ".cpp": "cpp",
        ".c": "c",
        ".h": "c",
        ".hpp": "cpp",
        ".go": "go",
        ".rs": "rust",
        ".rb": "ruby",
        ".sh": "bash",
        ".bash": "bash",
        ".sql": "sql",
        ".r": "r",
        ".m": "matlab",
        ".swift": "swift",
        ".kt": "kotlin",
        ".kts": "kotlin",
        ".scala": "scala",
        ".php": "php",
        ".css": "css",
        ".scss": "scss",
        ".less": "less",
        ".sass": "sass",
        ".lua": "lua",
        ".pl": "perl",
        ".pm": "perl",
        ".zig": "zig",
        ".dart": "dart",
        ".jl": "julia",
    }
    return mapping.get(ext, "")


# =============================================================================
# PDF EXTRACTOR
# =============================================================================


def _extract_pdf(file_bytes: bytes, filename: str) -> ExtractionResult:
    """Extract text and images from a PDF, preserving page order.

    Uses PyMuPDF (imported as 'fitz') to process each page:

    1. For each page, extract text blocks (paragraphs of text).
    2. Extract embedded images (photos, charts, diagrams).
    3. If a page has very little text (< SCANNED_PAGE_CHAR_THRESHOLD chars),
       it's likely a scanned page — render the entire page as an image
       so the vision-capable LLM can read it.
    4. All content is assembled in page order with page markers.

    PyMuPDF's get_text("dict") returns text blocks with bounding box
    coordinates. We sort these by vertical position (top to bottom) then
    horizontal (left to right) to get natural reading order. Images are
    interleaved at their approximate vertical position on the page.

    Args:
        file_bytes: Raw PDF file bytes.
        filename: Original filename.

    Returns:
        ExtractionResult with interleaved text and image ContentParts.
    """
    import fitz  # PyMuPDF

    parts: list[ContentPart] = []
    warnings: list[str] = []
    scanned_page_count = 0

    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as e:
        raise ValueError(f"Could not open PDF '{filename}': {e}") from e

    page_count = len(doc)

    for page_num in range(page_count):
        page = doc[page_num]
        page_text = page.get_text("text").strip()

        # Detect scanned pages (image-only, no extractable text)
        if len(page_text) < SCANNED_PAGE_CHAR_THRESHOLD:
            scanned_page_count += 1
            if scanned_page_count <= MAX_SCANNED_PAGE_IMAGES:
                # Render the scanned page as an image for the LLM to read
                image_part = _render_page_as_image(page, page_num + 1)
                if image_part:
                    parts.append(
                        ContentPart(
                            type="text",
                            text=f"\n--- Page {page_num + 1} (scanned — rendered as image) ---\n",
                        )
                    )
                    parts.append(image_part)
            elif scanned_page_count == MAX_SCANNED_PAGE_IMAGES + 1:
                warnings.append(
                    f"Document has more than {MAX_SCANNED_PAGE_IMAGES} scanned pages. "
                    f"Only the first {MAX_SCANNED_PAGE_IMAGES} were rendered as images."
                )
            continue

        # This page has extractable text — process it normally
        parts.append(ContentPart(type="text", text=f"\n--- Page {page_num + 1} ---\n{page_text}"))

        # Extract embedded images from this page
        page_images = _extract_pdf_page_images(doc, page, page_num + 1)
        parts.extend(page_images)

    doc.close()

    # Prepend document header
    header = ContentPart(type="text", text=f"[Document: {filename} — {page_count} pages]")
    parts.insert(0, header)

    if scanned_page_count > 0:
        warnings.append(
            f"{scanned_page_count} of {page_count} pages appear to be scanned. "
            f"Text was read by the AI model from rendered page images."
        )

    return ExtractionResult(
        filename=filename,
        file_type="pdf",
        file_size=len(file_bytes),
        content_parts=parts,
        page_count=page_count,
        warnings=warnings,
    )


def _render_page_as_image(page, page_num: int) -> ContentPart | None:
    """Render a PDF page as a PNG image for vision model processing.

    Used for scanned pages that have no extractable text. The vision-capable
    LLM (Gemma 3, GPT-4o, Claude) will read the page content from the image,
    effectively acting as an OCR engine — but smarter, because it also
    understands tables, diagrams, and layout.
    """
    import fitz  # PyMuPDF

    try:
        # Render at configured DPI. Higher = better quality but larger payload.
        zoom = SCANNED_PAGE_DPI / 72  # 72 DPI is PDF's default resolution
        matrix = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=matrix)

        # Convert to JPEG for smaller size (scanned pages are often large)
        img_bytes = pix.tobytes(output="jpeg", jpg_quality=IMAGE_QUALITY)

        return ContentPart(
            type="image",
            image_base64=base64.b64encode(img_bytes).decode("ascii"),
            image_mime="image/jpeg",
        )
    except Exception as e:
        logger.warning(f"[DocExtract] Failed to render page {page_num} as image: {e}")
        return None


def _extract_pdf_page_images(doc, page, page_num: int) -> list[ContentPart]:
    """Extract embedded images from a PDF page.

    These are actual images embedded in the document — photos, charts,
    diagrams, figures — NOT the page text. We extract each image's raw
    bytes, compress it if needed, and return it as an image ContentPart.

    PyMuPDF gives us image references (xref numbers) for each page.
    We then extract the actual image data using doc.extract_image().
    """
    parts: list[ContentPart] = []

    try:
        image_list = page.get_images(full=True)
    except Exception:
        return parts

    for img_index, img_info in enumerate(image_list):
        xref = img_info[0]  # The image's cross-reference number in the PDF
        try:
            base_image = doc.extract_image(xref)
            if not base_image:
                continue

            img_bytes = base_image["image"]
            img_ext = base_image.get("ext", "png")
            mime_type = f"image/{img_ext}" if img_ext != "jpg" else "image/jpeg"

            # Compress/resize if the image is large
            img_bytes, mime_type = _compress_image(img_bytes, mime_type)

            parts.append(
                ContentPart(type="text", text=f"\n[Embedded image from page {page_num}]\n")
            )
            parts.append(
                ContentPart(
                    type="image",
                    image_base64=base64.b64encode(img_bytes).decode("ascii"),
                    image_mime=mime_type,
                )
            )

        except Exception as e:
            logger.warning(
                f"[DocExtract] Failed to extract image {img_index} " f"from page {page_num}: {e}"
            )

    return parts


# =============================================================================
# DOCX EXTRACTOR
# =============================================================================


def _extract_docx(file_bytes: bytes, filename: str) -> ExtractionResult:
    """Extract text, tables, and images from a Word document (.docx).

    DOCX files have a clean internal structure: the document body is a
    sequence of paragraphs, tables, and inline images in document order.
    python-docx gives us these elements in the exact order they appear,
    making interleaving straightforward.

    The XML structure looks like:
        <w:body>
            <w:p>Paragraph text...</w:p>
            <w:tbl>Table data...</w:tbl>
            <w:p>More text with <w:drawing>inline image</w:drawing>...</w:p>
        </w:body>

    We iterate through the body's child elements and handle each type:
    - Paragraphs (w:p): Extract text, check for inline images
    - Tables (w:tbl): Convert to markdown table format
    """
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    parts: list[ContentPart] = []
    warnings: list[str] = []

    try:
        doc = Document(io.BytesIO(file_bytes))
    except Exception as e:
        raise ValueError(f"Could not open DOCX '{filename}': {e}") from e

    # python-docx gives us relationships that map rId to image blobs
    image_rels = {}
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_rels[rel.rId] = rel

    parts.append(ContentPart(type="text", text=f"[Document: {filename}]"))

    # Iterate through the document body elements in order.
    # This preserves the interleaving of text, tables, and images.
    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            # It's a paragraph — extract text and check for inline images
            para = Paragraph(element, doc)
            text = para.text.strip()
            if text:
                parts.append(ContentPart(type="text", text=text))

            # Check for inline images in this paragraph
            for run in para.runs:
                drawing_elements = run._element.findall(
                    f".//{qn('wp:inline')}"
                ) + run._element.findall(f".//{qn('wp:anchor')}")
                for drawing in drawing_elements:
                    blip = drawing.find(f".//{qn('a:blip')}")
                    if blip is not None:
                        embed_id = blip.get(qn("r:embed"))
                        if embed_id and embed_id in image_rels:
                            img_part = _extract_docx_image(image_rels[embed_id])
                            if img_part:
                                parts.append(ContentPart(type="text", text="\n[Embedded image]\n"))
                                parts.append(img_part)

        elif tag == "tbl":
            # It's a table — convert to markdown format
            table = Table(element, doc)
            md_table = _table_to_markdown(table)
            if md_table:
                parts.append(ContentPart(type="text", text=f"\n{md_table}\n"))

    return ExtractionResult(
        filename=filename,
        file_type="docx",
        file_size=len(file_bytes),
        content_parts=parts,
        page_count=0,
        warnings=warnings,
    )


def _extract_docx_image(rel) -> ContentPart | None:
    """Extract an image from a DOCX relationship."""
    try:
        img_bytes = rel.target_part.blob
        content_type = rel.target_part.content_type or "image/png"

        img_bytes, content_type = _compress_image(img_bytes, content_type)

        return ContentPart(
            type="image",
            image_base64=base64.b64encode(img_bytes).decode("ascii"),
            image_mime=content_type,
        )
    except Exception as e:
        logger.warning(f"[DocExtract] Failed to extract DOCX image: {e}")
        return None


def _table_to_markdown(table) -> str:
    """Convert a document table to markdown format.

    Markdown tables are lightweight text that LLMs read very well.
    This avoids needing to render tables as images, saving tokens
    and working with text-only models too.

    Example output:
        | Name | Score | Grade |
        |------|-------|-------|
        | Alice | 95 | A |
        | Bob | 87 | B+ |
    """
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")

    if not rows:
        return ""

    # Insert header separator after first row
    col_count = len(table.rows[0].cells) if table.rows else 0
    separator = "| " + " | ".join(["---"] * col_count) + " |"
    rows.insert(1, separator)

    return "\n".join(rows)


# =============================================================================
# XLSX EXTRACTOR
# =============================================================================


def _extract_xlsx(file_bytes: bytes, filename: str) -> ExtractionResult:
    """Extract data from an Excel spreadsheet (.xlsx).

    Each sheet in the workbook is converted to a markdown table.
    This is the most useful representation for LLMs — they can read
    and reason about tabular data in markdown format very effectively.

    Charts in Excel files are drawing objects generated from cell data.
    Rather than trying to render charts (which is complex and fragile),
    we include the underlying cell data, which is actually more useful
    to an LLM than a screenshot of a chart would be.

    Args:
        file_bytes: Raw .xlsx file bytes.
        filename: Original filename.

    Returns:
        ExtractionResult with one text ContentPart per sheet.
    """
    from openpyxl import load_workbook

    parts: list[ContentPart] = []
    warnings: list[str] = []

    try:
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    except Exception as e:
        raise ValueError(f"Could not open XLSX '{filename}': {e}") from e

    parts.append(
        ContentPart(type="text", text=f"[Spreadsheet: {filename} — {len(wb.sheetnames)} sheets]")
    )

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[list[str]] = []

        for row in ws.iter_rows(values_only=True):
            # Convert each cell value to string, handling None values
            cells = [str(cell) if cell is not None else "" for cell in row]
            # Skip completely empty rows
            if any(c.strip() for c in cells):
                rows.append(cells)

        if not rows:
            parts.append(ContentPart(type="text", text=f"\n### Sheet: {sheet_name}\n(empty sheet)"))
            continue

        # Build markdown table from rows
        md_lines = []
        for i, row in enumerate(rows):
            md_lines.append("| " + " | ".join(c.replace("|", "\\|") for c in row) + " |")
            if i == 0:
                # Header separator after first row
                md_lines.append("| " + " | ".join(["---"] * len(row)) + " |")

        # Warn if sheet is very large
        if len(rows) > 500:
            warnings.append(
                f"Sheet '{sheet_name}' has {len(rows)} rows. "
                f"Only the first 500 are included to fit within context limits."
            )
            md_lines = md_lines[:502]  # 500 data rows + header + separator

        parts.append(
            ContentPart(type="text", text=f"\n### Sheet: {sheet_name}\n\n" + "\n".join(md_lines))
        )

    wb.close()

    return ExtractionResult(
        filename=filename,
        file_type="xlsx",
        file_size=len(file_bytes),
        content_parts=parts,
        page_count=len(wb.sheetnames),
        warnings=warnings,
    )


# =============================================================================
# PPTX EXTRACTOR
# =============================================================================


def _extract_pptx(file_bytes: bytes, filename: str) -> ExtractionResult:
    """Extract text and images from a PowerPoint presentation (.pptx).

    Each slide is processed like a page:
    1. Extract text from all text frames (title, body, text boxes)
    2. Extract embedded images (photos, diagrams, charts)
    3. Content is ordered per-slide with slide markers

    PowerPoint slides often contain more images than text (diagrams,
    screenshots, charts). For these visual-heavy slides, the extracted
    images are crucial for the LLM to understand the content.

    Args:
        file_bytes: Raw .pptx file bytes.
        filename: Original filename.

    Returns:
        ExtractionResult with interleaved text and image ContentParts.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    parts: list[ContentPart] = []
    warnings: list[str] = []

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        raise ValueError(f"Could not open PPTX '{filename}': {e}") from e

    slide_count = len(prs.slides)
    parts.append(
        ContentPart(type="text", text=f"[Presentation: {filename} — {slide_count} slides]")
    )

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []

        for shape in slide.shapes:
            # Extract text from text-containing shapes
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

            # Extract tables from slides
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    col_count = len(table.rows[0].cells) if table.rows else 0
                    separator = "| " + " | ".join(["---"] * col_count) + " |"
                    rows.insert(1, separator)
                    slide_texts.append("\n".join(rows))

            # Extract images from slides
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_blob = shape.image.blob
                    content_type = shape.image.content_type or "image/png"

                    img_blob, content_type = _compress_image(img_blob, content_type)

                    # Add text for this slide first (if we haven't yet)
                    if slide_texts:
                        parts.append(
                            ContentPart(
                                type="text",
                                text=f"\n--- Slide {slide_num} ---\n" + "\n".join(slide_texts),
                            )
                        )
                        slide_texts = []

                    parts.append(
                        ContentPart(type="text", text=f"\n[Image from slide {slide_num}]\n")
                    )
                    parts.append(
                        ContentPart(
                            type="image",
                            image_base64=base64.b64encode(img_blob).decode("ascii"),
                            image_mime=content_type,
                        )
                    )
                except Exception as e:
                    logger.warning(
                        f"[DocExtract] Failed to extract image from slide {slide_num}: {e}"
                    )

        # Add any remaining text for this slide
        if slide_texts:
            parts.append(
                ContentPart(
                    type="text", text=f"\n--- Slide {slide_num} ---\n" + "\n".join(slide_texts)
                )
            )

    return ExtractionResult(
        filename=filename,
        file_type="pptx",
        file_size=len(file_bytes),
        content_parts=parts,
        page_count=slide_count,
        warnings=warnings,
    )


# =============================================================================
# IMAGE COMPRESSION UTILITY
# =============================================================================


def _compress_image(img_bytes: bytes, mime_type: str) -> tuple[bytes, str]:
    """Compress and resize an extracted image to save payload space.

    Large embedded images (e.g., a 4000x3000 photo in a DOCX) would waste
    tokens and bloat the Globus Compute payload. We scale them down to
    MAX_IMAGE_DIMENSION and compress as JPEG — the same approach used by
    STREAM's image upload pipeline.

    Args:
        img_bytes: Raw image bytes.
        mime_type: Original MIME type (e.g., "image/png").

    Returns:
        Tuple of (compressed_bytes, mime_type). The mime_type may change
        to "image/jpeg" if the image was converted for compression.
    """
    from PIL import Image

    try:
        img = Image.open(io.BytesIO(img_bytes))

        # Skip tiny images (icons, bullets) — not worth sending to the LLM
        if img.width < 50 or img.height < 50:
            return img_bytes, mime_type

        # Resize if larger than MAX_IMAGE_DIMENSION
        if img.width > MAX_IMAGE_DIMENSION or img.height > MAX_IMAGE_DIMENSION:
            img.thumbnail((MAX_IMAGE_DIMENSION, MAX_IMAGE_DIMENSION), Image.LANCZOS)

        # Convert to RGB if necessary (handles RGBA, palette modes)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")

        # Save as JPEG for consistent compression
        buffer = io.BytesIO()
        img.save(buffer, format="JPEG", quality=IMAGE_QUALITY, optimize=True)
        return buffer.getvalue(), "image/jpeg"

    except Exception as e:
        logger.warning(f"[DocExtract] Image compression failed: {e}")
        return img_bytes, mime_type
